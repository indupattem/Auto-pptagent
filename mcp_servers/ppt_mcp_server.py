from mcp.server.fastmcp import FastMCP
from pptx import Presentation
import os
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import requests

mcp = FastMCP("ppt-generator")

BASE_DIR = os.path.join(os.path.expanduser("~"), "AutoPPT_Files")
IMG_DIR = os.path.join(BASE_DIR, "images")

os.makedirs(BASE_DIR, exist_ok=True)
os.makedirs(IMG_DIR, exist_ok=True)

store = {}

# -------------------- CREATE PPT --------------------
@mcp.tool()
def create_presentation(filename: str) -> str:
    prs = Presentation()
    path = os.path.join(BASE_DIR, filename)
    store[filename] = (prs, path)
    return f"Created {filename}"


# -------------------- IMAGE DOWNLOAD --------------------

def download_image(query, filename):
    try:
        url = f"https://source.unsplash.com/800x500/?{query}"

        response = requests.get(url, timeout=5, allow_redirects=True)

        path = os.path.join(IMG_DIR, filename)

        with open(path, "wb") as f:
            f.write(response.content)

        return path

    except Exception as e:
        print("❌ Image download failed:", e)
        return None


# -------------------- ADD SLIDE --------------------
@mcp.tool()
def add_slide(filename: str, title: str) -> str:
    prs, _ = store[filename]
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # 🎨 DARK BACKGROUND (PROFESSIONAL)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 30, 60)

    # 🟦 TOP BAR
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(0.8),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 102, 204)
    rect.line.fill.background()

    # 🧠 TITLE STYLE
    title_shape = slide.shapes.title
    title_shape.text = title
    para = title_shape.text_frame.paragraphs[0]
    para.font.size = Pt(34)
    para.font.bold = True
    para.font.color.rgb = RGBColor(255, 255, 255)

    # 📌 CONTENT POSITION
    content = slide.placeholders[1]
    content.left = Inches(0.8)
    content.top = Inches(1.8)
    content.width = Inches(5.5)

    # 🖼️ BETTER IMAGE QUERY
    image_query = title + " " + "technology diagram"

    image_path = download_image(image_query, f"{title}.jpg")

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(6.2),
                Inches(2),
                width=Inches(3),
                height=Inches(2)
            )
        except:
            pass

    return f"Slide added: {title}"


# -------------------- WRITE CONTENT --------------------
@mcp.tool()
def write_content(filename: str, slide_index: int, points: list) -> str:
    prs, _ = store[filename]
    slide = prs.slides[slide_index]

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, p in enumerate(points):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.text = p
        para.font.size = Pt(18)
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.LEFT

    return f"Content added to slide {slide_index}"


# -------------------- SAVE PPT --------------------
@mcp.tool()
def save_presentation(filename: str) -> str:
    prs, path = store[filename]

    try:
        if os.path.exists(path):
            os.remove(path)

    except PermissionError:
        import time
        new_path = path.replace(".pptx", f"_{int(time.time())}.pptx")
        prs.save(new_path)
        return f"Saved at {new_path}"

    prs.save(path)
    return f"Saved at {path}"


# -------------------- RUN --------------------
if __name__ == "__main__":
    mcp.run()